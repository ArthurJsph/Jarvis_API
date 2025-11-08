"""
Manipulação de Apresentações PowerPoint
Jarvis CLI - Módulo de PowerPoint
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import json
import config
from rich.console import Console

console = Console()


class PowerPointManager:
    """Gerenciador de apresentações PowerPoint"""
    
    def __init__(self):
        self.encoding = config.DEFAULT_ENCODING
        self.default_author = config.DEFAULT_AUTHOR
    
    def create_presentation(self, filename: str, title: str = "", subtitle: str = "") -> bool:
        """
        Cria uma nova apresentação PowerPoint
        
        Args:
            filename: Nome do arquivo
            title: Título da apresentação
            subtitle: Subtítulo da apresentação
        """
        try:
            prs = Presentation()
            
            # Layout do slide de título
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Adiciona título e subtítulo
            if title:
                slide.shapes.title.text = title
            
            if subtitle and len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle
            
            # Salva a apresentação
            filepath = Path(config.DOCUMENTS_DIR) / filename
            if not filepath.suffix:
                filepath = filepath.with_suffix('.pptx')
            
            prs.save(filepath)
            
            console.print(f"✅ Apresentação PowerPoint criada: {filepath}", style="green")
            return True
            
        except Exception as e:
            console.print(f"❌ Erro ao criar apresentação: {str(e)}", style="red")
            return False
    
    def add_slide(self, filename: str, slide_title: str, content: str = "", 
                  layout_index: int = 1) -> bool:
        """
        Adiciona um slide a uma apresentação existente
        
        Args:
            filename: Nome do arquivo
            slide_title: Título do slide
            content: Conteúdo do slide
            layout_index: Índice do layout (0=título, 1=título e conteúdo, etc.)
        """
        try:
            filepath = Path(filename)
            if not filepath.exists():
                filepath = Path(config.DOCUMENTS_DIR) / filename
                
            if not filepath.exists():
                console.print(f"❌ Arquivo não encontrado: {filename}", style="red")
                return False
            
            prs = Presentation(filepath)
            
            # Adiciona novo slide
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            # Adiciona título
            if slide.shapes.title:
                slide.shapes.title.text = slide_title
            
            # Adiciona conteúdo se o layout suportar
            if content and len(slide.placeholders) > 1:
                slide.placeholders[1].text = content
            
            # Salva a apresentação
            prs.save(filepath)
            
            console.print(f"✅ Slide adicionado à apresentação: {filepath}", style="green")
            return True
            
        except Exception as e:
            console.print(f"❌ Erro ao adicionar slide: {str(e)}", style="red")
            return False
    
    def generate_from_excel(self, excel_file: str, ppt_file: str, 
                           title_column: str = None, content_column: str = None) -> bool:
        """
        Gera slides a partir de dados de uma planilha Excel
        
        Args:
            excel_file: Arquivo Excel com os dados
            ppt_file: Arquivo PowerPoint de saída
            title_column: Nome da coluna com títulos dos slides
            content_column: Nome da coluna com conteúdo dos slides
        """
        try:
            # Lê dados do Excel
            excel_path = Path(excel_file)
            if not excel_path.exists():
                excel_path = Path(config.DOCUMENTS_DIR) / excel_file
            
            df = pd.read_excel(excel_path)
            
            # Verifica se as colunas existem
            if title_column and title_column not in df.columns:
                console.print(f"❌ Coluna '{title_column}' não encontrada", style="red")
                return False
            
            if content_column and content_column not in df.columns:
                console.print(f"❌ Coluna '{content_column}' não encontrada", style="red")
                return False
            
            # Usa as primeiras duas colunas se não especificado
            if not title_column:
                title_column = df.columns[0]
            if not content_column and len(df.columns) > 1:
                content_column = df.columns[1]
            
            # Cria apresentação
            prs = Presentation()
            
            # Slide de título
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = f"Apresentação gerada de {excel_path.name}"
            
            # Cria slides a partir dos dados
            for index, row in df.iterrows():
                slide_title = str(row[title_column]) if title_column else f"Slide {index + 1}"
                slide_content = str(row[content_column]) if content_column else ""
                
                # Adiciona slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_title
                
                if slide_content and len(slide.placeholders) > 1:
                    slide.placeholders[1].text = slide_content
            
            # Salva apresentação
            ppt_path = Path(config.DOCUMENTS_DIR) / ppt_file
            prs.save(ppt_path)
            
            console.print(f"✅ Apresentação gerada com {len(df)} slides: {ppt_path}", style="green")
            return True
            
        except Exception as e:
            console.print(f"❌ Erro ao gerar apresentação do Excel: {str(e)}", style="red")
            return False
    
    def generate_from_json(self, json_file: str, ppt_file: str) -> bool:
        """
        Gera slides a partir de dados JSON
        
        Args:
            json_file: Arquivo JSON com os dados dos slides
            ppt_file: Arquivo PowerPoint de saída
        """
        try:
            # Lê dados do JSON
            json_path = Path(json_file)
            if not json_path.exists():
                json_path = Path(config.DOCUMENTS_DIR) / json_file
            
            with open(json_path, 'r', encoding=self.encoding) as f:
                data = json.load(f)
            
            # Cria apresentação
            prs = Presentation()
            
            # Slide de título principal
            if 'presentation' in data:
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                title_slide.shapes.title.text = data['presentation'].get('title', 'Apresentação')
                if 'subtitle' in data['presentation'] and len(title_slide.placeholders) > 1:
                    title_slide.placeholders[1].text = data['presentation']['subtitle']
            
            # Adiciona slides
            slides_data = data.get('slides', [])
            for slide_data in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # Título do slide
                if 'title' in slide_data:
                    slide.shapes.title.text = slide_data['title']
                
                # Conteúdo do slide
                if 'content' in slide_data and len(slide.placeholders) > 1:
                    if isinstance(slide_data['content'], list):
                        content = '\n'.join(slide_data['content'])
                    else:
                        content = str(slide_data['content'])
                    slide.placeholders[1].text = content
            
            # Salva apresentação
            ppt_path = Path(config.DOCUMENTS_DIR) / ppt_file
            prs.save(ppt_path)
            
            console.print(f"✅ Apresentação gerada com {len(slides_data)} slides: {ppt_path}", style="green")
            return True
            
        except Exception as e:
            console.print(f"❌ Erro ao gerar apresentação do JSON: {str(e)}", style="red")
            return False
    
    def extract_text(self, filename: str, output_file: Optional[str] = None) -> bool:
        """
        Extrai todo o texto de uma apresentação PowerPoint
        
        Args:
            filename: Nome do arquivo PowerPoint
            output_file: Arquivo de saída (opcional)
        """
        try:
            filepath = Path(filename)
            if not filepath.exists():
                filepath = Path(config.DOCUMENTS_DIR) / filename
                
            if not filepath.exists():
                console.print(f"❌ Arquivo não encontrado: {filename}", style="red")
                return False
            
            prs = Presentation(filepath)
            
            extracted_text = []
            
            for i, slide in enumerate(prs.slides, 1):
                extracted_text.append(f"=== SLIDE {i} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            extracted_text.append(text)
                
                # Extrai texto das notas do slide
                if slide.notes_slide.notes_text_frame.text.strip():
                    extracted_text.append("--- NOTAS ---")
                    extracted_text.append(slide.notes_slide.notes_text_frame.text.strip())
                
                extracted_text.append("")  # Linha em branco entre slides
            
            # Junta todo o texto
            full_text = '\n'.join(extracted_text)
            
            # Salva em arquivo se especificado
            if not output_file:
                output_file = filepath.stem + "_extracted.txt"
            
            output_path = Path(config.DOCUMENTS_DIR) / output_file
            
            with open(output_path, 'w', encoding=self.encoding) as f:
                f.write(full_text)
            
            console.print(f"✅ Texto extraído para: {output_path}", style="green")
            console.print(f"Slides processados: {len(prs.slides)}", style="cyan")
            return True
            
        except Exception as e:
            console.print(f"❌ Erro ao extrair texto: {str(e)}", style="red")
            return False
    
    def add_chart_slide(self, filename: str, slide_title: str, chart_data: Dict[str, List]) -> bool:
        """
        Adiciona um slide com gráfico simples
        
        Args:
            filename: Nome do arquivo
            slide_title: Título do slide
            chart_data: Dados do gráfico {"labels": [...], "values": [...]}
        """
        try:
            filepath = Path(filename)
            if not filepath.exists():
                filepath = Path(config.DOCUMENTS_DIR) / filename
                
            if not filepath.exists():
                console.print(f"❌ Arquivo não encontrado: {filename}", style="red")
                return False
            
            prs = Presentation(filepath)
            
            # Adiciona slide em branco
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
            
            # Adiciona título
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            
            # Cria tabela simples como "gráfico"
            labels = chart_data.get('labels', [])
            values = chart_data.get('values', [])
            
            if labels and values:
                rows = len(labels) + 1  # +1 para cabeçalho
                cols = 2
                
                table = slide.shapes.add_table(rows, cols, Inches(2), Inches(2), Inches(6), Inches(4)).table
                
                # Cabeçalho
                table.cell(0, 0).text = "Item"
                table.cell(0, 1).text = "Valor"
                
                # Dados
                for i, (label, value) in enumerate(zip(labels, values)):
                    table.cell(i + 1, 0).text = str(label)
                    table.cell(i + 1, 1).text = str(value)
            
            # Salva apresentação
            prs.save(filepath)
            
            console.print(f"✅ Slide com dados adicionado: {filepath}", style="green")
            return True
            
        except Exception as e:
            console.print(f"❌ Erro ao adicionar slide com gráfico: {str(e)}", style="red")
            return False
    
    def get_presentation_info(self, filename: str) -> Dict[str, Any]:
        """
        Obtém informações sobre uma apresentação PowerPoint
        
        Args:
            filename: Nome do arquivo
        """
        try:
            filepath = Path(filename)
            if not filepath.exists():
                filepath = Path(config.DOCUMENTS_DIR) / filename
                
            if not filepath.exists():
                console.print(f"❌ Arquivo não encontrado: {filename}", style="red")
                return {}
            
            prs = Presentation(filepath)
            
            # Conta elementos
            total_slides = len(prs.slides)
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            
            # Conta texto
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            
            full_text = ' '.join(all_text)
            words = len(full_text.split())
            
            # Informações do arquivo
            file_size = filepath.stat().st_size
            
            info = {
                "filename": filepath.name,
                "size": f"{file_size / 1024:.2f} KB",
                "slides": total_slides,
                "shapes": total_shapes,
                "words": words,
                "characters": len(full_text),
                "layouts_available": len(prs.slide_layouts)
            }
            
            return info
            
        except Exception as e:
            console.print(f"❌ Erro ao obter info da apresentação: {str(e)}", style="red")
            return {}
    
    def create_template_json(self, filename: str = "presentation_template.json") -> bool:
        """
        Cria um arquivo JSON template para gerar apresentações
        
        Args:
            filename: Nome do arquivo template
        """
        try:
            template = {
                "presentation": {
                    "title": "Minha Apresentação",
                    "subtitle": "Criada com Jarvis CLI"
                },
                "slides": [
                    {
                        "title": "Introdução",
                        "content": [
                            "• Primeiro ponto",
                            "• Segundo ponto",
                            "• Terceiro ponto"
                        ]
                    },
                    {
                        "title": "Desenvolvimento",
                        "content": "Conteúdo do slide de desenvolvimento..."
                    },
                    {
                        "title": "Conclusão",
                        "content": [
                            "• Resumo dos pontos principais",
                            "• Próximos passos",
                            "• Perguntas?"
                        ]
                    }
                ]
            }
            
            filepath = Path(config.DOCUMENTS_DIR) / filename
            
            with open(filepath, 'w', encoding=self.encoding) as f:
                json.dump(template, f, indent=2, ensure_ascii=False)
            
            console.print(f"✅ Template JSON criado: {filepath}", style="green")
            console.print("Use: jarvis gerar ppt do json [arquivo.json] [saida.pptx]", style="cyan")
            return True
            
        except Exception as e:
            console.print(f"❌ Erro ao criar template JSON: {str(e)}", style="red")
            return False